from pptx import Presentation
from pptx.util import Inches

pr = Presentation()

# Slide 1: Title Slide
slide = pr.slides.add_slide(pr.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Epidemiology, Treatment, and Control of Tuberculosis"
subtitle.text = "Comprehensive Teaching Module for MBBS Students"

# Slide 2: Learning Objectives
slide = pr.slides.add_slide(pr.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "Learning Objectives"

tf = body_shape.text_frame
tf.text = "By the end of this module, students should be able to:"

p = tf.add_paragraph()
p.text = "1. Describe the global epidemiology of tuberculosis"
p.level = 1

p = tf.add_paragraph()
p.text = "2. Explain the transmission, pathogenesis, and risk factors for TB"
p.level = 1

p = tf.add_paragraph()
p.text = "3. Identify clinical manifestations and diagnostic approaches"
p.level = 1

p = tf.add_paragraph()
p.text = "4. Outline treatment regimens for drug-susceptible and drug-resistant TB"
p.level = 1

p = tf.add_paragraph()
p.text = "5. Discuss preventive measures and national TB control programs"
p.level = 1

# Slide 3: Introduction
slide = pr.slides.add_slide(pr.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "Introduction"

tf = body_shape.text_frame
tf.text = "Tuberculosis (TB) is an infectious disease caused by Mycobacterium tuberculosis bacteria."

p = tf.add_paragraph()
p.text = "Primarily affects the lungs (pulmonary TB) but can involve other organs (extrapulmonary TB)"
p.level = 0

p = tf.add_paragraph()
p.text = "TB is preventable and curable but remains a major global health challenge"
p.level = 0

# Slide 4: Epidemiology
slide = pr.slides.add_slide(pr.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "Epidemiology"

tf = body_shape.text_frame
tf.text = "Global Burden (WHO 2024):"

p = tf.add_paragraph()
p.text = "- 10 million new TB cases annually"
p.level = 1

p = tf.add_paragraph()
p.text = "- 1.5 million TB-related deaths per year"
p.level = 1

p = tf.add_paragraph()
p.text = "- Top infectious disease killer worldwide"
p.level = 1

p = tf.add_paragraph()
p.text = "- Major contributor to antimicrobial resistance"
p.level = 1

p = tf.add_paragraph()
p.text = "High burden countries: Bangladesh, China, India, Indonesia, Nigeria, Pakistan, Philippines, South Africa"
p.level = 0

p = tf.add_paragraph()
p.text = "India: Leads with ~2.7 million cases annually"
p.level = 0

# Slide 5: Transmission and Pathogenesis
slide = pr.slides.add_slide(pr.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "Transmission and Pathogenesis"

tf = body_shape.text_frame
tf.text = "Transmission:"

p = tf.add_paragraph()
p.text = "Through air: Droplet nuclei from coughing, sneezing, or speaking by active TB patients"
p.level = 0

p = tf.add_paragraph()
p.text = "Highly contagious in close, prolonged contact"
p.level = 0

p = tf.add_paragraph()
p.text = "Pathogenesis:"
p.level = 0

p = tf.add_paragraph()
p.text = "Inhalation → macrophages ingest bacteria → primary infection in lungs"
p.level = 1

p = tf.add_paragraph()
p.text = "Latent TB: Bacteria dormant, can reactivate with immunosuppression"
p.level = 1

p = tf.add_paragraph()
p.text = "Active TB: Progressive infection, tissue destruction"
p.level = 1

# Slide 6: Risk Factors
slide = pr.slides.add_slide(pr.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "Risk Factors"

tf = body_shape.text_frame
tf.text = "Host Factors:"

p = tf.add_paragraph()
p.text = "- HIV/AIDS (25% of HIV deaths due to TB)"
p.level = 1

p = tf.add_paragraph()
p.text = "- Malnutrition and underweight"
p.level = 1

p = tf.add_paragraph()
p.text = "- Diabetes mellitus"
p.level = 1

p = tf.add_paragraph()
p.text = "- Tobacco use and smoking"
p.level = 1

p = tf.add_paragraph()
p.text = "- Alcohol use disorder"
p.level = 1

p = tf.add_paragraph()
p.text = "Social and Environmental Factors:"
p.level = 0

p = tf.add_paragraph()
p.text = "- Overcrowding, poor ventilation"
p.level = 1

p = tf.add_paragraph()
p.text = "- Poverty, homelessness"
p.level = 1

p = tf.add_paragraph()
p.text = "- Occupational exposure (healthcare workers)"
p.level = 1

# Slide 7: Clinical Manifestations
slide = pr.slides.add_slide(pr.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "Clinical Manifestations"

tf = body_shape.text_frame
tf.text = "Pulmonary TB Symptoms:"

p = tf.add_paragraph()
p.text = "- Persistent cough (>2 weeks)"
p.level = 1

p = tf.add_paragraph()
p.text = "- Chest pain"
p.level = 1

p = tf.add_paragraph()
p.text = "- Hemoptysis (coughing blood)"
p.level = 1

p = tf.add_paragraph()
p.text = "- Fatigue and weakness"
p.level = 1

p = tf.add_paragraph()
p.text = "- Unexplained weight loss"
p.level = 1

p = tf.add_paragraph()
p.text = "- Fever and night sweats"
p.level = 1

p = tf.add_paragraph()
p.text = "Extrapulmonary TB: Lymph nodes, bones, kidneys, etc."
p.level = 0

# Slide 8: Diagnosis
slide = pr.slides.add_slide(pr.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "Diagnosis"

tf = body_shape.text_frame
tf.text = "Suspected cases: Clinical history + chest X-ray"

p = tf.add_paragraph()
p.text = "Microbiological Confirmation:"
p.level = 0

p = tf.add_paragraph()
p.text = "- Sputum smear microscopy (smear-positive rate)"
p.level = 1

p = tf.add_paragraph()
p.text = "- Rapid molecular tests (GeneXpert MTB/RIF recommended initial test)"
p.level = 1

p = tf.add_paragraph()
p.text = "- Culture for definitive diagnosis"
p.level = 1

p = tf.add_paragraph()
p.text = "- Chest radiography, CT scans for complications"
p.level = 1

p = tf.add_paragraph()
p.text = "For Latent TB Infection:"
p.level = 0

p = tf.add_paragraph()
p.text = "- TST (Tuberculin Skin Test)"
p.level = 1

p = tf.add_paragraph()
p.text = "- IGRA (Interferon-Gamma Release Assays)"
p.level = 1

# Slide 9: Treatment (DS-TB)
slide = pr.slides.add_slide(pr.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "Treatment: Drug-Susceptible TB"

tf = body_shape.text_frame
tf.text = "Standard Regimen (Directly Observed Treatment - DOTS):"

p = tf.add_paragraph()
p.text = "- Intensive Phase: 2 months Isoniazid (H), Rifampicin (R), Pyrazinamide (Z), Ethambutol (E)"
p.level = 1

p = tf.add_paragraph()
p.text = "- Continuation Phase: 4 months Isoniazid (H), Rifampicin (R)"
p.level = 1

p = tf.add_paragraph()
p.text = "Duration: 6 months total"
p.level = 0

p = tf.add_paragraph()
p.text = "Key Principles:"
p.level = 0

p = tf.add_paragraph()
p.text = "- Supervised treatment to ensure adherence"
p.level = 1

p = tf.add_paragraph()
p.text = "- Fixed-dose combinations to simplify"
p.level = 1

p = tf.add_paragraph()
p.text = "- Monitor for side effects and treatment response"
p.level = 1

# Slide 10: Treatment (DR-TB)
slide = pr.slides.add_slide(pr.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "Treatment: Drug-Resistant TB"

tf = body_shape.text_frame
tf.text = "Multidrug-Resistant TB (MDR-TB): Resistant to H and R"

p = tf.add_paragraph()
p.text = "- Intensive Phase: 6-8 months with injectable (Kanamycin/Streptomycin) + 4 second-line drugs"
p.level = 1

p = tf.add_paragraph()
p.text = "- Continuation Phase: 12-18 months with oral agents"
p.level = 1

p = tf.add_paragraph()
p.text = "- Total duration: 18-24 months"
p.level = 1

p = tf.add_paragraph()
p.text = "Extensively Drug-Resistant TB (XDR-TB): Resists fluoroquinolones and injectables too"
p.level = 0

p = tf.add_paragraph()
p.text = "Challenges: Longer treatment, higher toxicity, lower cure rates (<50%)"
p.level = 0

p = tf.add_paragraph()
p.text = "Access to new drugs: Bedaquiline, Delamanid"
p.level = 0

# Slide 11: Prevention and Control
slide = pr.slides.add_slide(pr.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "Prevention and Control"

tf = body_shape.text_frame
tf.text = "Primary Prevention:"

p = tf.add_paragraph()
p.text = "- BCG vaccine for infants in high-burden countries"
p.level = 1

p = tf.add_paragraph()
p.text = "- Control of risk factors (HIV, diabetes management, tobacco cessation)"
p.level = 1

p = tf.add_paragraph()
p.text = "Secondary Prevention:"
p.level = 0

p = tf.add_paragraph()
p.text = "- TB Preventive Treatment (TPT) for LTBI: 3-6 months INH or 3-6HP"
p.level = 1

p = tf.add_paragraph()
p.text = "Tertiary Prevention:"
p.level = 0

p = tf.add_paragraph()
p.text = "- Effective treatment to prevent transmission"
p.level = 1

p = tf.add_paragraph()
p.text = "- Contact tracing and screening"
p.level = 1

# Slide 12: Special Populations at Risk
slide = pr.slides.add_slide(pr.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "TB in Special Populations"

tf = body_shape.text_frame
tf.text = "High-risk groups requiring special attention:"

p = tf.add_paragraph()
p.text = "• Pediatric TB: 1.2M cases, 30-50% extrapulmonary"
p.level = 1

p = tf.add_paragraph()
p.text = "• TB-HIV Co-infection: 25% of HIV deaths due to TB"
p.level = 1

p = tf.add_paragraph()
p.text = "• TB in Diabetes: 2-3X increased risk"
p.level = 1

p = tf.add_paragraph()
p.text = "• TB in Pregnancy: Requires modified regimens"
p.level = 1

p = tf.add_paragraph()
p.text = "• Healthcare Workers: 2-5X higher risk"
p.level = 1

# Slide 12: TB Control in India - NTEP
slide = pr.slides.add_slide(pr.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "TB Control in India: National TB Elimination Programme (NTEP)"

tf = body_shape.text_frame
tf.text = "Evolution: NTP (1962) → RNTCP (1992) → NTEP (2020)"

p = tf.add_paragraph()
p.text = "Goals: Eliminate TB by 2030 (<1 case per million)"
p.level = 0

p = tf.add_paragraph()
p.text = "90-90-90 Targets: Detect, Treat, Cure 90% cases"
p.level = 1

p = tf.add_paragraph()
p.text = "Progress (2023): 2.3M cases notified, 84% cure rate"
p.level = 1

# Slide 13: NTEP Strategies
slide = pr.slides.add_slide(pr.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "NTEP Key Strategies"

tf = body_shape.text_frame
tf.text = "Five Pillars of NTEP:"

p = tf.add_paragraph()
p.text = "1. Early Diagnosis: CBNAAT as initial test"
p.level = 1

p = tf.add_paragraph()
p.text = "2. Treatment Support: Free drugs, DOT, nutrition"
p.level = 1

p = tf.add_paragraph()
p.text = "3. Prevention: TPT, BCG, risk factor control"
p.level = 1

p = tf.add_paragraph()
p.text = "4. Surveillance: Nikshay 2.0 portal"
p.level = 1

p = tf.add_paragraph()
p.text = "5. Public-Private Partnerships"
p.level = 1

# Slide 14: NTEP Progress & Challenges
slide = pr.slides.add_slide(pr.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "NTEP Progress & Challenges"

tf = body_shape.text_frame
tf.text = "Successes:"

p = tf.add_paragraph()
p.text = "• 20% reduction in TB incidence (2015-2023)"
p.level = 1

p = tf.add_paragraph()
p.text = "• Private sector notifications increased"
p.level = 1

p = tf.add_paragraph()
p.text = "• Innovations: TrueNat, 99DOTS, AI diagnostics"
p.level = 1

p = tf.add_paragraph()
p.text = "Challenges:"
p.level = 0

p = tf.add_paragraph()
p.text = "• Detection gap: 0.5M undiagnosed cases"
p.level = 1

p = tf.add_paragraph()
p.text = "• Socioeconomic barriers, stigma"
p.level = 1

p = tf.add_paragraph()
p.text = "• Health system strengthening needed"
p.level = 1

# Slide 13: Summary
slide = pr.slides.add_slide(pr.slide_layouts[5])
shapes = slide.shapes
shapes.title.text = "Summary"

tb_summary = [
    "TB is a leading cause of death from an infectious disease",
    "Early diagnosis and prompt treatment are crucial",
    "DOTS strategy effective for curing patients and reducing transmission",
    "Prevention requires vaccination, risk factor reduction, and LTBI treatment",
    "Strong health systems and community engagement essential for control"
]

text_box = shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
tf = text_box.text_frame
for point in tb_summary:
    p = tf.add_paragraph()
    p.text = point

# Slide 14: References
slide = pr.slides.add_slide(pr.slide_layouts[5])
shapes = slide.shapes
shapes.title.text = "References"

ref_text = [
    "WHO Global Tuberculosis Report 2023",
    "CDC Tuberculosis Fact Sheet",
    "National TB Program Guidelines, India",
    "Textbook of Preventive and Social Medicine (MBD Muraleedharan)"
]

text_box = shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
tf = text_box.text_frame
tf.text = "Sources:"
for ref in ref_text:
    p = tf.add_paragraph()
    p.text = ref

pr.save('presentations/TB_Module_Presentation.pptx')

print("Presentation saved as TB_Module_Presentation.pptx")
